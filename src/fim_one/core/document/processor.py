"""Vision-aware document processing.

Provides text extraction and PDF page rendering for both traditional
text-based and vision-model-based document understanding.
"""

from __future__ import annotations

import asyncio
import base64
import json
import logging
import os
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Environment configuration
# ---------------------------------------------------------------------------


def _get_doc_processing_mode() -> str:
    """Return the configured document processing mode (auto | vision | text)."""
    return os.environ.get("DOCUMENT_PROCESSING_MODE", "auto")


def _get_doc_vision_dpi() -> int:
    """Return the configured DPI for PDF page rendering."""
    return int(os.environ.get("DOCUMENT_VISION_DPI", "150"))


def _get_doc_vision_max_pages() -> int:
    """Return the maximum number of PDF pages to render as images."""
    return int(os.environ.get("DOCUMENT_VISION_MAX_PAGES", "20"))


# ---------------------------------------------------------------------------
# Data types
# ---------------------------------------------------------------------------

IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg"}


@dataclass
class DocumentResult:
    """Result of processing a document through the vision-aware pipeline."""

    text: str | None
    page_images: list[str] = field(default_factory=list)  # base64 data URLs
    mode_used: str = "text"  # "vision" | "text"
    page_count: int = 0


# ---------------------------------------------------------------------------
# Main processor
# ---------------------------------------------------------------------------


class DocumentProcessor:
    """Vision-aware document processing.

    Combines traditional text extraction with optional PDF page rendering
    for vision-capable LLMs.
    """

    @staticmethod
    async def extract_text(file_path: Path) -> str | None:
        """Extract text content from a file.

        Supports plain text, JSON, CSV, PDF (via pdfplumber), and
        Office documents (via markitdown). Images return ``None``.

        Args:
            file_path: Path to the file to extract text from.

        Returns:
            Extracted text content, a fallback message when optional
            dependencies are missing, or ``None`` for unsupported types.
        """
        return await asyncio.to_thread(_extract_text_sync, file_path)

    @staticmethod
    async def render_pdf_pages(
        file_path: Path,
        dpi: int = 150,
        max_pages: int = 20,
    ) -> list[bytes]:
        """Render PDF pages as PNG images using PyMuPDF (fitz).

        Args:
            file_path: Path to the PDF file.
            dpi: Rendering resolution (default 150).
            max_pages: Maximum number of pages to render (default 20).

        Returns:
            List of PNG image bytes, one per rendered page.

        Raises:
            ImportError: If PyMuPDF is not installed.
        """

        def _render() -> list[bytes]:
            import fitz  # type: ignore[import-untyped]

            doc = fitz.open(str(file_path))
            images: list[bytes] = []
            try:
                for i, page in enumerate(doc):
                    if i >= max_pages:
                        break
                    mat = fitz.Matrix(dpi / 72, dpi / 72)
                    pix = page.get_pixmap(matrix=mat)
                    images.append(pix.tobytes("png"))
            finally:
                doc.close()
            return images

        return await asyncio.to_thread(_render)

    @staticmethod
    async def process_document(
        file_path: Path,
        mode: str = "auto",
        supports_vision: bool = False,
        max_pages: int | None = None,
        dpi: int | None = None,
    ) -> DocumentResult:
        """Process a document with a vision-aware strategy.

        Args:
            file_path: Path to the document file.
            mode: Processing mode — ``"vision"``, ``"text"``, or ``"auto"``
                (default). In auto mode, vision is used when the model
                supports it and the document is a PDF.
            supports_vision: Whether the active model supports vision.
            max_pages: Maximum PDF pages to render (falls back to env var).
            dpi: Rendering DPI (falls back to env var).

        Returns:
            A :class:`DocumentResult` with extracted text and optional
            rendered page images.
        """
        if max_pages is None:
            max_pages = _get_doc_vision_max_pages()
        if dpi is None:
            dpi = _get_doc_vision_dpi()

        suffix = file_path.suffix.lower()

        # Always extract text (needed for both modes)
        text = await DocumentProcessor.extract_text(file_path)

        # Determine effective mode
        if mode == "auto":
            effective = (
                "vision" if supports_vision and suffix == ".pdf" else "text"
            )
        else:
            effective = mode

        page_images: list[str] = []
        page_count = 0

        if effective == "vision" and suffix == ".pdf":
            try:
                raw_images = await DocumentProcessor.render_pdf_pages(
                    file_path, dpi, max_pages
                )
                page_count = len(raw_images)
                page_images = [
                    f"data:image/png;base64,{base64.b64encode(img).decode('ascii')}"
                    for img in raw_images
                ]
            except ImportError:
                logger.warning(
                    "PyMuPDF not installed, falling back to text-only "
                    "document processing. Install with: uv add PyMuPDF"
                )
                effective = "text"
            except Exception:
                logger.warning(
                    "PDF rendering failed, falling back to text",
                    exc_info=True,
                )
                effective = "text"

        if not page_images:
            effective = "text"

        return DocumentResult(
            text=text,
            page_images=page_images,
            mode_used=effective,
            page_count=page_count,
        )

    @staticmethod
    async def extract_with_images(file_path: Path) -> tuple[str | None, list[bytes]]:
        """Extract text with embedded images from a document.

        For DOCX/PPTX files, returns text with ``[Figure N]`` markers and
        the corresponding image bytes.  For all other formats, delegates to
        :func:`_extract_text_sync` and returns an empty image list.

        Args:
            file_path: Path to the document file.

        Returns:
            A tuple of ``(text_with_figure_markers, list_of_image_bytes)``.
        """
        return await asyncio.to_thread(_extract_with_images_sync, file_path)

    @staticmethod
    async def get_or_create_cached_pages(
        file_path: Path,
        dpi: int | None = None,
        max_pages: int | None = None,
    ) -> list[str]:
        """Return cached PDF page images, rendering them if not yet cached.

        Page images are stored as individual PNG files under a ``.pages/``
        directory next to the original file for fast retrieval during chat.

        Args:
            file_path: Path to the PDF file.
            dpi: Rendering DPI (falls back to env var).
            max_pages: Maximum pages to render (falls back to env var).

        Returns:
            List of base64 data URL strings for each page.
        """
        if dpi is None:
            dpi = _get_doc_vision_dpi()
        if max_pages is None:
            max_pages = _get_doc_vision_max_pages()

        pages_dir = file_path.parent / ".pages" / file_path.stem
        pages_dir.mkdir(parents=True, exist_ok=True)

        # Check for existing cached pages
        existing = sorted(pages_dir.glob("page_*.png"))
        if existing:
            data_urls: list[str] = []
            for pg in existing[:max_pages]:
                raw = await asyncio.to_thread(pg.read_bytes)
                b64 = base64.b64encode(raw).decode("ascii")
                data_urls.append(f"data:image/png;base64,{b64}")
            return data_urls

        # Render and cache
        try:
            raw_images = await DocumentProcessor.render_pdf_pages(
                file_path, dpi, max_pages
            )
        except ImportError:
            logger.warning("PyMuPDF not installed, cannot render PDF pages")
            return []
        except Exception:
            logger.warning("PDF rendering failed", exc_info=True)
            return []

        data_urls = []
        for i, img_bytes in enumerate(raw_images):
            page_file = pages_dir / f"page_{i:04d}.png"
            await asyncio.to_thread(page_file.write_bytes, img_bytes)
            b64 = base64.b64encode(img_bytes).decode("ascii")
            data_urls.append(f"data:image/png;base64,{b64}")

        return data_urls


# ---------------------------------------------------------------------------
# Sync text extraction (delegated from the upload flow)
# ---------------------------------------------------------------------------


def _extract_text_sync(file_path: Path) -> str | None:
    """Synchronous text extraction implementation.

    This is the single source of truth for text extraction logic.
    Both the async :meth:`DocumentProcessor.extract_text` and the
    upload-time ``_extract_content`` call this function.
    """
    suffix = file_path.suffix.lower()

    # Images have no extractable text content
    if suffix in IMAGE_EXTENSIONS:
        return None

    # Plain text family
    if suffix in {".txt", ".md", ".py", ".js", ".html", ".htm"}:
        return file_path.read_text(encoding="utf-8", errors="replace")

    # JSON -- parse and pretty-print
    if suffix == ".json":
        try:
            data = json.loads(
                file_path.read_text(encoding="utf-8", errors="replace")
            )
            return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            return file_path.read_text(encoding="utf-8", errors="replace")

    # CSV -- raw text
    if suffix == ".csv":
        return file_path.read_text(encoding="utf-8", errors="replace")

    # PDF -- requires pdfplumber (optional)
    if suffix == ".pdf":
        try:
            import pdfplumber
        except ImportError:
            return "[PDF content extraction requires pdfplumber]"
        pages_text: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages_text.append(text)
        return "\n".join(pages_text) if pages_text else None

    # Office documents (DOCX, DOC, XLSX, XLS, PPTX, PPT) -- requires markitdown
    if suffix in {".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt"}:
        try:
            from markitdown import MarkItDown
        except ImportError:
            return (
                f"[{suffix.upper().lstrip('.')} content extraction "
                f"requires markitdown]"
            )
        converter = MarkItDown()
        result = converter.convert(str(file_path))
        content = result.text_content or ""
        return content if content.strip() else None

    return None


# ---------------------------------------------------------------------------
# DOCX/PPTX embedded image extraction with positional references
# ---------------------------------------------------------------------------


def _extract_docx_with_images(file_path: Path) -> tuple[str, list[bytes]]:
    """Extract text with positional image markers from DOCX.

    Iterates paragraphs and tables in document order.  When a paragraph
    contains embedded images (``<a:blip>`` elements), a ``[Figure N]``
    marker is inserted at the image's position and the raw image bytes
    are collected.

    Returns:
        A tuple of ``(text_with_markers, list_of_image_bytes)``.
    """
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(str(file_path))
    parts: list[str] = []
    images: list[bytes] = []
    fig_num = 0

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]  # strip namespace

        if tag == "p":  # paragraph
            para = Paragraph(element, doc)

            # Detect embedded images via <a:blip> elements
            blips = element.findall(
                ".//"
                "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
            )

            para_text = para.text.strip()

            if blips:
                for blip in blips:
                    embed = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/"
                        "2006/relationships}embed"
                    )
                    if embed:
                        try:
                            rel = doc.part.rels[embed]
                            image_bytes = rel.target_part.blob
                            fig_num += 1
                            images.append(image_bytes)
                            parts.append(f"\n[Figure {fig_num}]\n")
                        except (KeyError, Exception):
                            pass

            if para_text:
                parts.append(para_text)

        elif tag == "tbl":  # table
            table = Table(element, doc)
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            if rows:
                parts.append("\n".join(rows))

    return "\n\n".join(parts), images


def _extract_pptx_with_images(file_path: Path) -> tuple[str, list[bytes]]:
    """Extract text with positional image markers from PPTX.

    Iterates slides in order.  Picture shapes produce ``[Figure N]``
    markers; grouped shapes are also inspected for nested images.

    Returns:
        A tuple of ``(text_with_markers, list_of_image_bytes)``.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(file_path))
    parts: list[str] = []
    images: list[bytes] = []
    fig_num = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = [f"--- Slide {slide_idx} ---"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_texts.append(text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    fig_num += 1
                    images.append(image_bytes)
                    slide_texts.append(f"[Figure {fig_num}]")
                except Exception:
                    pass

            # Inspect grouped shapes for nested images
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if hasattr(child, "image"):
                        try:
                            fig_num += 1
                            images.append(child.image.blob)
                            slide_texts.append(f"[Figure {fig_num}]")
                        except Exception:
                            pass

        parts.append("\n".join(slide_texts))

    return "\n\n".join(parts), images


def _extract_with_images_sync(file_path: Path) -> tuple[str | None, list[bytes]]:
    """Extract text with embedded images from documents.

    For DOCX/PPTX files, returns text with ``[Figure N]`` markers and
    the corresponding image bytes.  For all other formats, delegates to
    :func:`_extract_text_sync` and returns an empty image list.

    Returns:
        A tuple of ``(text_with_figure_markers, list_of_image_bytes)``.
    """
    suffix = file_path.suffix.lower()

    if suffix in (".docx", ".doc"):
        try:
            return _extract_docx_with_images(file_path)
        except Exception:
            logger.warning(
                "DOCX image extraction failed, falling back to text-only",
                exc_info=True,
            )
            return _extract_text_sync(file_path), []

    if suffix in (".pptx", ".ppt"):
        try:
            return _extract_pptx_with_images(file_path)
        except Exception:
            logger.warning(
                "PPTX image extraction failed, falling back to text-only",
                exc_info=True,
            )
            return _extract_text_sync(file_path), []

    # All other formats: text only, no embedded images
    return _extract_text_sync(file_path), []
